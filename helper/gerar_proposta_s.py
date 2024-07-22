import pptx
from datetime import datetime
from clientes.models import Cliente
from ucs.models import ModuloSolar, Projeto, UC, Endereco, Categoria, TipoUC
from .format_functions import numberFormatter
from .utils import obter_dados, obter_dimensionamento, calcular_payback
from .gerar_utils import update_text_of_textbox, update_text_of_table, update_image, update_chart, formatar_payback
import io
import os


#class Cliente():
    
#    def __init__(self, nome, cidade, estado):
#        self.nome = nome
#        self.cidade = cidade
#        self.estado = estado



# class Inversor():

#     def __init__(self, marca, modelo, potencia, correnteMax, garantiaAno, 
#                  codInmetro, categoria="Inversor"):
#         self.categoria = categoria
#         self.marca = marca
#         self.modelo = modelo
#         self.potencia = potencia
#         self.correnteMax = correnteMax
#         self.garantiaAno = garantiaAno
#         self.codInmetro = codInmetro


# class ModuloSolar():

#     def __init__(self, marca, modelo, potencia, area, peso, garantiaAno,
#                  codInmetro):
#         self.marca = marca
#         self.modelo = modelo
#         self.potencia = potencia
#         self.area = area
#         self.peso = peso
#         self.garantiaAno = garantiaAno
#         self.codInmetro = codInmetro


# class Projeto():

#     def __init__(self, consumoTotal, qtdeModulos, producaoMedia, qtdeInv):
#         self.consumoTotal = consumoTotal
#         self.qtdeModulos = qtdeModulos
#         self.producaoMedia = producaoMedia
#         self.qtdeInv = qtdeInv



# Get the current directory of the script
current_dir = os.path.dirname(os.path.abspath(__file__))
# Construct the path to the template file
template_folder = os.path.join(current_dir, 'templates')
template_path = os.path.join(template_folder, 'Template_Proposta.pptx')
PRESENTATION = pptx.Presentation(template_path)









#if __name__ == '__main__':
#    cl = Cliente.objects.all()
#    print(cl)



def gerar_PPTX(uc: UC, projeto:Projeto):
    '''Função que gera automaticamente a proposta simples normal em PPTX e salva no computador'''
    ### PÁGINA 1
    # NOME E DATA NO TEXT_BOX 1
    modulo = projeto.modulo
    cliente = uc.cliente
    investimento = int(projeto.valorProposta)
    nome = cliente.nome.upper() # parametro
    data = datetime.today().strftime('%d/%m/%Y')
    nome_data_text = f"""{nome}
{data}"""


    update_text_of_textbox(PRESENTATION, 1, 1, nome_data_text)

    ### PÁGINA 2
    
    # página 2 alteração no ID: 4 (<<invoumicro>>)
    invoumicro = "Inversor" # hard CODED
    invoumicro_text = f'{invoumicro}:' #precisa inserir como parametro
    update_text_of_textbox(PRESENTATION, 2, 4, invoumicro_text)


    # tabela ID 1 => CELL IDs 
    # ( 2: <<marcasModulo>>, 4: <<potMod>>, 6: <<Qtde_Mod>>, 
    #   8: <<pesoMod>>, 12: <<grantMod>>, 14: <<inmetroMod>>)
    # Pot_inv no table 2 cell id 4
    # Qtde_inv no table 2 cell id 6

    # Tabela 1 => ID 2
    marcasModulo = ['Ja Solar', 'TSUN', 'Sunova'] # SET IT AS FIXED FOR NOW
    marcasModulo_text = ""
    for idx, marca in enumerate(marcasModulo):
        marcasModulo_text += marca + '/' if idx + 1 != len(marcasModulo) else marca

    update_text_of_table(PRESENTATION, 2, 1, 2, marcasModulo_text)

    # Tabela 1 => ID 4
    potMod = int(modulo.potencia)
    potMod_text = f"{potMod} Wp"

    update_text_of_table(PRESENTATION, 2, 1, 4, potMod_text)

    # Tabela 1 => ID 6
    qtde_mod = projeto.qtdeModulos
    qtde_text = f"{qtde_mod}"
    update_text_of_table(PRESENTATION, 2, 1, 6, qtde_text)

    # Tabela 1 => ID 8
    peso_Mod = numberFormatter(modulo.peso)
    peso_Mod_text = f"{peso_Mod} kg"

    update_text_of_table(PRESENTATION, 2, 1, 8, peso_Mod_text)

    # Tabela 1 => ID 12
    garantiaMod = modulo.garantiaAno
    garantiaMod_text = f"{garantiaMod} anos"

    update_text_of_table(PRESENTATION, 2, 1, 12, garantiaMod_text)

    # Tabela 1 => ID 14
    inmetroMod = modulo.codInmetro

    update_text_of_table(PRESENTATION, 2, 1, 14, inmetroMod)

    #  tabela ID 2 => CELL IDs ( 2: <<marcasInv>>, 4: <<Pot_Inv>>,
    #  6: <<qtdeInv>>, 12: <<grantInv>>, 14: <<inmetroInv>> )

    # Tabela 2 => ID 2
    marcasInv = ['Growatt', 'Solis', 'Sofar'] # SET IT AS FIXED FOR NOW
    marcasInv_text = ""
    for idx, marca in enumerate(marcasInv):
        marcasInv_text += marca + '/' if idx + 1 != len(marcasInv) else marca
    
    update_text_of_table(PRESENTATION, 2, 2, 2, marcasInv_text)

    # Tabela 2 => ID 4
    pot_inv, _ = obter_dimensionamento(projeto.qtdeModulos) # SET IT AS 1 FOR NOW
    pot_text = ""
    for idx, inv in enumerate(pot_inv, 1):
        pot_text += str(inv) + ' kW'
        if idx != len(pot_inv):
            pot_text += ' + '

    update_text_of_table(PRESENTATION, 2, 2, 4, pot_text)

    # Tabela 2 => ID 6
    qtde_inversores_text = str(projeto.qtdeInv)
    update_text_of_table(PRESENTATION, 2, 2, 6, qtde_inversores_text)

    # Tabela 2 => ID 12
    #print(f"Printing before error {inversor}")
    garantiaInv = "10" # Hard Coded
    garantiaInv_text = f"{garantiaInv} anos"
    update_text_of_table(PRESENTATION, 2, 2, 12, garantiaInv_text)

    # Tabela 2 => ID 14
    inmetroInv = "009133/2019" # Hard Coded
    update_text_of_table(PRESENTATION, 2, 2, 14, inmetroInv)

    # (alteração de imagem) (DONE)
    # IMAGE ID 11
    image_id = 11
    # remover shape cinze que fica na frente da imagem se for microinversor
    if invoumicro == 'Microinversor':
        #print('Removing additional shape')
        shapes = PRESENTATION.slides[2 - 1].shapes
        shapes.element.remove(shapes[12-1].element)
        img_path = os.path.join(template_folder, 'microinversor.png')
    else:
        img_path = os.path.join(template_folder, 'inversor.png')

    update_image(PRESENTATION, 2, image_id, img_path)

    ### PÁGINA 3
    # POT_MOD, CIDADE E ESTADO ESTÃO NO TEXT_BOX 3 (check)
    # PROD_MED NO TEXT_BOX 6  (check)
    # GRÁFICO EM IMAGEM A SER COLOCADO

    pot_mod = int(qtde_mod) * int(modulo.potencia) / 1000
    pot_mod_formatado = "{:,.2f}".format(pot_mod).replace('.', ',')
    cidade = uc.endereco.cidade
    estado = uc.endereco.estado
    pot_cidade_estado_text = f"Determinamos que o sistema solar fotovoltaico apontado neste relatório, com potencia de {pot_mod_formatado} kWp, é apto a produzir em média na cidade de {cidade}/{estado}."

    update_text_of_textbox(PRESENTATION, 3, 3, pot_cidade_estado_text)

    prod_media = numberFormatter(projeto.producaoMedia)
    prod_med_text = f"""Produção média
    mensal de energia:
    {prod_media} kWh"""

    update_text_of_textbox(PRESENTATION, 3, 6, prod_med_text)

    ## COLOCAR CHART AQUI 
    new_data = obter_dados(int(projeto.consumoTotal)) # parametro
    update_chart(PRESENTATION, 3, new_data)

    ### PÁGINA 4
    # PROPOSTA NO TEXT_BOX 26
    # PAYBACK NO TEXT_BOX 25

    proposta = investimento # parametro
    proposta_formatada = "{:,.2f}".format(proposta).replace('.', ',').replace(',', '.', 1)

    proposta_text = f" {proposta_formatada}"
    update_text_of_textbox(PRESENTATION, 4, 26, proposta_text)

    anos, meses = calcular_payback(float(projeto.producaoMedia), investimento) # parametro

    payback_formatado = formatar_payback(anos, meses)

    payback_text = f"{payback_formatado}"
    update_text_of_textbox(PRESENTATION, 4, 25, payback_text)

    ### PÁGINA 5
    # alteração no ID 3 e no ID 7 (<<invoumicro>>, <<garantiamod>>, <<garantiaInv>>)
    text_id3 = f"""Instalação dos painéis;
Instalação do {invoumicro.lower()};
Instalação de proteções elétricas;
Passagem de condutores;
Memorial descritivo;
ART’s de projeto e execução;
Treinamento do sistema de monitoramento (quando houver)."""
    
    update_text_of_textbox(PRESENTATION, 5, 3, text_id3)


    text_id7 = f"""Painel fotovoltaico: {garantiaMod} anos e garantia de fornecimento de 25 anos;
{invoumicro}: {garantiaInv} anos;
Estrutura: 12 anos;
Instalação: 1 ano (incluindo uma manutenção)."""

    update_text_of_textbox(PRESENTATION, 5, 7, text_id7)


    # Save the presentation to an in-memory file
    pptx_io = io.BytesIO()
    PRESENTATION.save(pptx_io)
    pptx_io.seek(0)

    return pptx_io



if __name__ == "__main__":
    cliente = Cliente(nome="Pedro Ribeiro", cpf="05215456500")

    endereco = Endereco(CEP="49045080", prefixo_local="Rua", rua="Armando Barros",
                        num_logradouro="81", bairro="Luzia", cidade="Aracaju", estado="SE")


    categoria = Categoria(nomeCategoria="Bifásico")
    tipouc = TipoUC(nomeTipo="Usina")
    uc = UC(num_UC="3/99992-8", cliente=cliente, endereco=endereco,
            categoria=categoria, tipoUC=tipouc, consumo=650, tensaoNominal="127/220", resideoucomercial="Residencial")

    sunova550 = ModuloSolar(marca='Sunova', modelo='SS-550-72-MDH', potencia=550,
                        area=2.68, peso=27.5, garantiaAno=15, 
                        codInmetro='003361/2024')

    projeto = Projeto(uc=uc, modulo=sunova550, consumoTotal=700, qtdeModulos=10, producaoMedia=750, 
                    qtdeInv=1)

    gerar_PPTX(uc=uc, projeto=projeto)
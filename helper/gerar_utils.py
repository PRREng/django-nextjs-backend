from pptx.chart.data import CategoryChartData

def list_text_boxes(presentation, slide_num):
    slide = presentation.slides[slide_num - 1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape.text)
    return text_boxes




def update_text_of_textbox(presentation, slide_num, text_box_id, new_text):
    slide = presentation.slides[(slide_num - 1)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Preserve formatting of the first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_color = font.color
                #print(f" PÁGINA {slide_num} , TEXT BOX {text_box_id}")
                #print(f" TYPE OF FONT COLOR: {type(font_color.theme_color)}")
                # Clear existing text and apply new text with preserved formatting
                text_frame.clear() # Clears all text and formatting
                new_run = text_frame.paragraphs[0].add_run() # New run in first paragraph
                new_run.text = new_text
                # Reapply formatting
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline


                try:
                    new_run.font.color.rgb = font_color.rgb 
                except:
                    print("RGB doesnt work")
                
                    try:
                        new_run.font.color.theme_color = font_color.theme_color
                    except:
                        print('Theme Color doesn\'t work')
                        print(f'Página {slide_num}, TextBox {text_box_id}')

                return



def update_text_of_table(presentation, slide_num, table_id, cell_id, new_text):
    
    slide = presentation.slides[slide_num - 1]
    count_table_id = 0
    for shape in slide.shapes:
        if shape.has_table and shape.table:
            #print("ITERATING TABLE")
            count_table_id += 1
            if count_table_id == table_id:
                #print(f"MATCHED ID: {table_id}")
                for idx, cell in enumerate(shape.table.iter_cells(), 1):
                    if idx == cell_id:
                        # change the text
                        text_frame = cell.text_frame
                        first_paragraph = text_frame.paragraphs[0]
                        first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                        # Preserve formatting of the first run
                        font = first_run.font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_italic = font.italic
                        font_underline = font.underline
                        font_color = font.color
                        #print(f" PÁGINA {slide_num} , TABLE ID {table_id}, CELL ID {cell_id}")
                        # print(f" TYPE OF FONT COLOR: {type(font_color.theme_color)}")
                        # Clear existing text and apply new text with preserved formatting
                        text_frame.clear() # Clears all text and formatting
                        new_run = text_frame.paragraphs[0].add_run() # New run in first paragraph
                        new_run.text = new_text
                        # Reapply formatting
                        new_run.font.name = font_name
                        new_run.font.size = font_size
                        new_run.font.bold = font_bold
                        new_run.font.italic = font_italic
                        new_run.font.underline = font_underline
                        try:
                            new_run.font.color.rgb = font_color.rgb 
                        except:
                            print("RGB doesnt work")
                        
                            try:
                                new_run.font.color.theme_color = font_color.theme_color
                            except:
                                print('Theme Color doesn\'t work')
                                print(f'Página {slide_num}, Tabela {table_id}, Celula {cell_id}')

                        return

    return





def update_chart(presentation, slide_num, new_data):
    slide = presentation.slides[slide_num - 1]

    # Access the chart (assuming it's a bar chart)
    for shape in slide.shapes:
        if shape.has_chart and shape.chart:
            chart = shape.chart

    chart_data = CategoryChartData()
    for label, value in new_data:
        chart_data.add_category(label="")
        chart_data.add_series(label, value)

    chart.replace_data(chart_data)

    return


def update_image(presentation, slide_num, shape_num, img_path):
    if not img_path:
        return None
    slide = presentation.slides[slide_num - 1]
    img = slide.shapes[shape_num - 1]
    try:
        imgPic = img._pic
    except AttributeError:
        raise AttributeError(
            f"Error for slide: {slide_num}, shape: {shape_num}, path: {img_path}")
    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
    imgPart = slide.part.related_part(imgRID)

    with open(img_path, 'rb') as f:
        rImgBlob = f.read()

    # replace
    imgPart._blob = rImgBlob
    #print(f"UPDATED THE IMAGE ID {shape_num}")







def list_tables(presentation, slide_num):
    slide = presentation.slides[slide_num - 1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_table and shape.table:
            text_boxes.append(shape.table)
    return text_boxes



def list_charts(presentation, slide_num):
    slide = presentation.slides[slide_num - 1]
    charts = []
    for shape in slide.shapes:
        if shape.has_chart and shape.chart:
            charts.append(shape.chart)
    return charts


def formatar_payback(anos, meses):
    if meses == 12:
        anos += 1
    payback_formatado = ""
    if anos > 1:
        payback_formatado += str(anos) + " anos"
    elif anos == 1:
        payback_formatado += str(anos) + " ano"

    if meses != 12:
        payback_formatado += " e "
        if meses > 1:
            payback_formatado += str(meses) + " meses"
        elif meses == 1:
            payback_formatado += str(meses) + " mês"
    return payback_formatado
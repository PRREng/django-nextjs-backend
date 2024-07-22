import pptx
from datetime import datetime
from ucs.models import Projeto, UC
from .gerar_utils import update_text_of_textbox, update_text_of_table, update_chart, update_image, formatar_payback
from .utils import obter_dados, obter_media_gerada, obter_dimensionamento, \
    calcular_payback, obter_econ_mensal, obter_econ_5anos, obter_econ_25anos
from .format_functions import numberFormatter, cpfFormatter
import io
import os

#==============================INIT===================
# Get the current directory of the script
current_dir = os.path.dirname(os.path.abspath(__file__))
# Construct the path to the template file
template_folder = os.path.join(current_dir, 'templates')
template_path = os.path.join(template_folder, 'Template_Grande.pptx')
PRESENTATION = pptx.Presentation(template_path)


#=========================================================

def gerarGrandePptx(uc: UC, projeto: Projeto):

    # -------------------------------------Página 1 (2-COMPLETED)----------------------------------------------
    # Text box, ID 1
    cliente = uc.cliente
    modulo = projeto.modulo
    investimento = int(projeto.valorProposta)

    nome = cliente.nome
    data = datetime.today().strftime('%d/%m/%Y')
    nome_data_text = f"""{nome}
{data}"""
    
    update_text_of_textbox(PRESENTATION, 1, 1, nome_data_text)


    # -------------------------------------Página 2 (1-COMPLETED)------------------------------------------------

    # Tabela ID 1 Cells (IDs => 4, 6, 8, 10, 12)
    # Cell 4
    nomecliente_text = f"{nome}."

    update_text_of_table(PRESENTATION, 2, 1, 4, nomecliente_text)

    # Cell 6
    cpf = cpfFormatter(cliente.cpf)
    cpfcnpj_text = f"{cpf}."

    update_text_of_table(PRESENTATION, 2, 1, 6, cpfcnpj_text)

    # Cell 8
    endereco = uc.endereco
    bairro = endereco.bairro
    cidade = endereco.cidade
    estado = endereco.estado
    bairro_cidade_text = f"{bairro}/{cidade}/{estado}."

    update_text_of_table(PRESENTATION, 2, 1, 8, bairro_cidade_text)

    # Cell 10
    rua = endereco.prefixo_local + ' ' + endereco.rua
    numero = endereco.num_logradouro
    rua_num_text = f"{rua}, {numero}."

    update_text_of_table(PRESENTATION, 2, 1, 10, rua_num_text)

    # Cell 12
    cep = endereco.CEP
    cep_text = f"{cep}."

    update_text_of_table(PRESENTATION, 2, 1, 12, cep_text)

    # Tabela ID 2 Cells (IDs => 4, 6, 8, 9, 10, 11, 12, 13, 14, 16, 18)
    # Cell 4
    prodMod = numberFormatter(projeto.qtdeModulos * modulo.potencia / 1000)
    prodmod_text = f"{prodMod} kWp."

    update_text_of_table(PRESENTATION, 2, 2, 4, prodmod_text)

    # Cell 6
    potMod = int(modulo.potencia)
    potmod_text = f"{potMod} Wp."

    update_text_of_table(PRESENTATION, 2, 2, 6, potmod_text)

    # Cell 8
    qtdeMod = projeto.qtdeModulos
    qtdemod_text = f"{qtdeMod}."

    update_text_of_table(PRESENTATION, 2, 2, 8, qtdemod_text)

    # Cell 9
    invoumicro = "Inversor"
    pot_invoumicro_text = f"Potência do {invoumicro.lower()}:"

    update_text_of_table(PRESENTATION, 2, 2, 9, pot_invoumicro_text)

    # Cell 10
    pot_inv, _ = obter_dimensionamento(projeto.qtdeModulos)
    pot_text = ""
    for idx, inv in enumerate(pot_inv, 1):
        pot_text += str(inv) + ' kW'
        if idx != len(pot_inv):
            pot_text += ' + '

    update_text_of_table(PRESENTATION, 2, 2, 10, pot_text)

    # Cell 11
    marcainvoumicro_text = f"Marca do {invoumicro.lower()}:"

    update_text_of_table(PRESENTATION, 2, 2, 11, marcainvoumicro_text)

    # Cell 12
    marcasInv = ['Growatt', 'Solis', 'Sofar'] # SET IT AS FIXED FOR NOW
    marcasInv_text = ""
    for idx, marca in enumerate(marcasInv):
        marcasInv_text += marca + '/' if idx + 1 != len(marcasInv) else marca

    update_text_of_table(PRESENTATION, 2, 2, 12, marcasInv_text)

    # Cell 13
    qtdeinvmicro_text = f"Quantidade de {invoumicro.lower()}:"

    update_text_of_table(PRESENTATION, 2, 2, 13, qtdeinvmicro_text)

    # Cell 14
    qtdeInv = projeto.qtdeInv
    qtdeinv_text = f"{qtdeInv}."

    update_text_of_table(PRESENTATION, 2, 2, 14, qtdeinv_text)

    # Cell 16
    areasist = numberFormatter(modulo.area * projeto.qtdeModulos)
    areasist_text = f"{areasist} m²."

    update_text_of_table(PRESENTATION, 2, 2, 16, areasist_text)

    # Cell 18
    pesoarea = numberFormatter(modulo.peso / modulo.area)
    pesoarea_text = f"{pesoarea} kg/m²."

    update_text_of_table(PRESENTATION, 2, 2, 18, pesoarea_text)


    # Tabela ID 3 Cells (IDs => 4, 6, 8, 10)
    # Cell 4
    resideoucomer = uc.resideoucomercial
    resideoucomer_text = f"{resideoucomer}."

    update_text_of_table(PRESENTATION, 2, 3, 4, resideoucomer_text)

    # Cell 6
    tensaonom = uc.tensaoNominal
    tensaonom_text = f"{tensaonom} V."

    update_text_of_table(PRESENTATION, 2, 3, 6, tensaonom_text)

    # Cell 8
    codUc = uc.num_UC
    coduc_text = f"{codUc}."

    update_text_of_table(PRESENTATION, 2, 3, 8, coduc_text)

    # Cell 10
    monobiftrif = uc.categoria
    monobiftrif_text = f"{monobiftrif}."

    update_text_of_table(PRESENTATION, 2, 3, 10, monobiftrif_text)


    #--------------------------------------Página 3 (1-COMPLETED)------------------------------------------------
    # Text box ID 3
    prodMod = numberFormatter(projeto.qtdeModulos * modulo.potencia / 1000)
    cidade = "Aracaju"
    estado = "SE"
    prod_cidade_text = f"""Determinamos que o sistema solar fotovoltaico apontado neste relatório, com potencia de {prodMod} kWp, \
é apto a produzir em média na cidade de {cidade}/{estado}."""
    
    update_text_of_textbox(PRESENTATION, 3, 3, prod_cidade_text)

    # Text box ID 6
    prodenergia = obter_media_gerada(projeto.consumoTotal)
    prod_ener_text = f"""Produção média
mensal de energia:
{prodenergia} kWh"""
    
    update_text_of_textbox(PRESENTATION, 3, 6, prod_ener_text)


    # Alteração de Chart
    new_data = obter_dados(projeto.consumoTotal)
    update_chart(PRESENTATION, 3, new_data)


    #---------------------------------------Página 4 (1-COMPLETE)---------------------------------------------
    # Text box ID 4
    #invoumicro = "Microinversor"
    invoumicro_text = f"""{invoumicro}:"""

    update_text_of_textbox(PRESENTATION, 4, 4, invoumicro_text)

    # possui 2 tabelas (DONE)
    # Tabela ID 1 Cells (IDs => 2, 4, 6, 10, 12) (DONE)
    # Cell 2
    marcaMod = modulo.marca
    marcamod_text = f"{marcaMod}"

    update_text_of_table(PRESENTATION, 4, 1, 2, marcamod_text)

    # Cell 4
    potmod_text = f"{potMod} Wp"

    update_text_of_table(PRESENTATION, 4, 1, 4, potmod_text)

    # Cell 6
    pesoMod = numberFormatter(modulo.peso)
    pesomod_text = f"{pesoMod} kg"

    update_text_of_table(PRESENTATION, 4, 1, 6, pesomod_text)

    # Cell 10
    garantiaMod = modulo.garantiaAno
    garantiaMod_text = f"{garantiaMod} anos"

    update_text_of_table(PRESENTATION, 4, 1, 10, garantiaMod_text)

    # Cell 12
    inmetroMod = modulo.codInmetro
    inmetromod_text = f"{inmetroMod}"

    update_text_of_table(PRESENTATION, 4, 1, 12, inmetromod_text)


    # Tabela ID 2 Cells (IDs => 2, 4, 10, 12) (DONE)
    # Cell 2
    # marcasInv = ['Growatt', 'Solis', 'Sofar'] # SET IT AS FIXED FOR NOW
    # marcasInv_text = ""
    # for idx, marca in enumerate(marcasInv):
    #     marcasInv_text += marca + '/' if idx + 1 != len(marcasInv) else marca

    update_text_of_table(PRESENTATION, 4, 2, 2, marcasInv_text)

    # Cell 4

    update_text_of_table(PRESENTATION, 4, 2, 4, pot_text)

    # Cell 10
    garantiaInv = "10" # Hard Coded
    garantiainv_text = f"{garantiaInv} anos"

    update_text_of_table(PRESENTATION, 4, 2, 10, garantiainv_text)

    # Cell 12
    inmetroInv = "009133/2019" # Hard Coded
    inmetroinv_text = f"{inmetroInv}"

    update_text_of_table(PRESENTATION, 4, 2, 12, inmetroinv_text)


    # há alteração de imagem (inversor) (Done)
    # IMAGE Shape ID => 11
    image_id = 11
    # remover shape cinze que fica na frente da imagem se for microinversor
    if invoumicro == 'Microinversor':
        # print('Removing additional shape')
        shapes = PRESENTATION.slides[4 - 1].shapes
        shapes.element.remove(shapes[12-1].element)
        img_path = os.path.join(template_folder, 'microinversor.png')
    else:
        img_path = os.path.join(template_folder, 'inversor.png')

    update_image(PRESENTATION, 4, image_id, img_path)



    #--------------------------------------Página 5 (1-COMPLETE)-------------------------------------------------
    # possui 1 tabela (DONE)
    # Tabela 1 IDs (Cells => 6)
    # Cell 6
    # pesoarea = numberFormatter(12.5) # formatar
    pesoarea_text = f"{pesoarea} kg/m²"

    update_text_of_table(PRESENTATION, 5, 1, 6, pesoarea_text)




    #---------------------------------------Página 6(1-COMPLETE)-----------------------------------------------
    # Text box ID 9
    valor_proposta = numberFormatter(investimento) # formatar (TODO)
    valor_proposta_text = f"""R$ {valor_proposta}"""

    update_text_of_textbox(PRESENTATION, 6, 9, valor_proposta_text)




    #----------------------------------------Página 7 (1-COMPLETE)-----------------------------------------------
    # Text box ID 7
    anos, meses = calcular_payback(float(projeto.producaoMedia), investimento) # parametro

    payback_extenso = formatar_payback(anos, meses)
    payback_text = f"""PAYBACK: {payback_extenso}"""

    update_text_of_textbox(PRESENTATION, 7, 7, payback_text)

    # Text box ID 8
    econMensal = numberFormatter(obter_econ_mensal(projeto.consumoTotal))
    econMensal_text = f"""R$ {econMensal}"""

    update_text_of_textbox(PRESENTATION, 7, 8, econMensal_text)

    # Text box ID 9
    econ5anos = numberFormatter(obter_econ_5anos(projeto.consumoTotal))
    econ5anos_text = f"""R$ {econ5anos}"""

    update_text_of_textbox(PRESENTATION, 7, 9, econ5anos_text)


    # Text box ID 10
    econ25anos = numberFormatter(obter_econ_25anos(projeto.consumoTotal))
    econ25anos_text = f"""R$ {econ25anos}"""

    update_text_of_textbox(PRESENTATION, 7, 10, econ25anos_text)




    #---------------------------------------Página 8 (1-COMPLETE)------------------------------------------------
    # Text box ID 3
    if projeto.isPlural():
        dosinvoumicro = f"dos {invoumicro.lower()}s"
    else:
        dosinvoumicro = f"do {invoumicro.lower()}"

    dosinvoumicro_text = f"""Instalação dos painéis;
Instalação {dosinvoumicro};
Instalação de proteções elétricas;
Passagem de condutores;
Memorial descritivo;
ART’s de projeto e execução;
Treinamento do sistema de monitoramento (quando houver)."""
    
    update_text_of_textbox(PRESENTATION, 8, 3, dosinvoumicro_text)

    # Text box ID 7
    garantias_text = f"""Painel fotovoltaico: {garantiaMod} anos e garantia de fornecimento de 25 anos;
{invoumicro}: {garantiaInv} anos;
Estrutura: 12 anos;
Instalação: 1 ano (incluindo uma manutenção)."""

    update_text_of_textbox(PRESENTATION, 8, 7, garantias_text)

    #-------------------------------------------DONE-------------------------------------------------


    # save the presentation to an in-memory file
    pptx_io = io.BytesIO()
    PRESENTATION.save(pptx_io)
    pptx_io.seek(0)

    return pptx_io










#======================================   Init  ====================================================== 



if __name__ == '__main__':
    gerarGrandePptx(PRESENTATION)

    # slide_num = 4
    # slide = PRESENTATION.slides[slide_num - 1]

    # # Picture shape IDs => 2, 5, 10, 11

    # for idx, shape in enumerate(slide.shapes, 1):
    #     print(f"Shape ID {idx}: {shape}")



    # print(f"PÁGINA: {slide_num}")
    # for idx, table in enumerate(tables):
    #     print(f"Tabela {idx + 1}")

    #     for i, cell in enumerate(table.iter_cells(), 1):
    #         print(f"Cell ID {i}: {cell.text}")